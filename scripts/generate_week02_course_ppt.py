from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "week02" / "week02_data_bss_course.pptx"

COLORS = {
    "ink": RGBColor(34, 37, 43),
    "muted": RGBColor(91, 99, 112),
    "bg": RGBColor(248, 249, 251),
    "line": RGBColor(216, 221, 230),
    "blue": RGBColor(28, 83, 160),
    "green": RGBColor(42, 126, 78),
    "orange": RGBColor(184, 95, 34),
    "teal": RGBColor(0, 116, 130),
    "red": RGBColor(164, 68, 64),
    "code_bg": RGBColor(32, 36, 44),
    "code_fg": RGBColor(238, 241, 246),
    "white": RGBColor(255, 255, 255),
}


def apply_font(run, size=22, bold=False, color=None, name="Microsoft YaHei"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]


def set_speaker_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text
    return slide


def set_background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["bg"]


def add_text(slide, x, y, w, h, body, size=22, bold=False, color=None, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = body
    apply_font(r, size=size, bold=bold, color=color, name=font_name)
    return box


def add_title(slide, heading, subtitle=None, lesson=None):
    if lesson:
        add_text(slide, 0.7, 0.28, 2.4, 0.3, lesson, size=13, bold=True, color=COLORS["blue"])
    add_text(slide, 0.7, 0.55, 12.0, 0.58, heading, size=30, bold=True)
    if subtitle:
        add_text(slide, 0.72, 1.14, 11.9, 0.35, subtitle, size=14, color=COLORS["muted"])
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(1.48), Inches(11.95), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.fill.background()


def add_bullet_slide(prs, heading, items, subtitle=None, lesson=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, heading, subtitle, lesson)
    box = slide.shapes.add_textbox(Inches(0.95), Inches(1.85), Inches(11.45), Inches(5.15))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(9)
        r = p.add_run()
        r.text = item
        apply_font(r, size=20 if len(item) < 34 else 17)
    return slide


def add_code_slide(prs, heading, code_text, subtitle=None, lesson=None, font_size=13):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, heading, subtitle, lesson)
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.82),
        Inches(1.75),
        Inches(11.72),
        Inches(5.08),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["code_bg"]
    shape.line.color.rgb = COLORS["code_bg"]
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.16)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = code_text
    apply_font(r, size=font_size, color=COLORS["code_fg"], name="Consolas")
    return slide


def add_code_annotated_slide(prs, heading, code_text, annotations, subtitle=None, lesson=None, font_size=12):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, heading, subtitle, lesson)
    code_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(1.75),
        Inches(6.6),
        Inches(5.08),
    )
    code_shape.fill.solid()
    code_shape.fill.fore_color.rgb = COLORS["code_bg"]
    code_shape.line.color.rgb = COLORS["code_bg"]
    tf = code_shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.14)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = code_text
    apply_font(r, size=font_size, color=COLORS["code_fg"], name="Consolas")

    box = slide.shapes.add_textbox(Inches(7.5), Inches(1.85), Inches(5.15), Inches(4.95))
    note_tf = box.text_frame
    note_tf.clear()
    note_tf.word_wrap = True
    for i, item in enumerate(annotations):
        para = note_tf.paragraphs[0] if i == 0 else note_tf.add_paragraph()
        para.space_after = Pt(8)
        run = para.add_run()
        run.text = item
        apply_font(run, size=14)
    return slide


def add_section_slide(prs, lesson, heading, goal):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_text(slide, 0.85, 1.1, 2.3, 0.35, lesson, size=15, bold=True, color=COLORS["blue"])
    add_text(slide, 0.85, 1.72, 11.4, 0.78, heading, size=36, bold=True)
    add_text(slide, 0.9, 2.86, 10.9, 0.95, goal, size=23, color=COLORS["muted"])
    add_text(slide, 0.9, 6.55, 11.4, 0.28, "本节结束时，学生要能用自己的话解释核心概念。", size=13, color=COLORS["muted"])
    return slide


def add_flow_slide(prs, heading, steps, subtitle=None, lesson=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, heading, subtitle, lesson)
    n = len(steps)
    widths = {3: 3.1, 4: 2.5, 5: 2.18, 6: 1.72}
    width = widths.get(n, 2.0)
    total_w = n * width + (n - 1) * 0.25
    x = (13.333 - total_w) / 2
    gap = 0.25
    for i, (name, body, color) in enumerate(steps):
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.35), Inches(width), Inches(1.75))
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLORS["white"]
        rect.line.color.rgb = COLORS["line"]
        add_text(slide, x + 0.16, 2.53, width - 0.32, 0.32, name, size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.18, 2.96, width - 0.36, 0.82, body, size=13, align=PP_ALIGN.CENTER)
        x += width + gap
        if i < n - 1:
            add_text(slide, x - 0.05, 3.0, 0.25, 0.3, "→", size=20, bold=True, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    return slide


def add_table_slide(prs, heading, subtitle, lesson, headers, rows, col_widths, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, heading, subtitle, lesson)
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.8), Inches(1.9), Inches(11.75), Inches(4.25))
    table = table_shape.table
    for idx, width in enumerate(col_widths):
        table.columns[idx].width = Inches(width)
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["blue"]
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                apply_font(r, size=13, bold=True, color=COLORS["white"])
    for row_idx, row in enumerate(rows, start=1):
        for col, value in enumerate(row):
            cell = table.cell(row_idx, col)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["white"]
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    apply_font(r, size=12)
    add_text(slide, 0.82, 6.55, 11.7, 0.35, footer, size=13, color=COLORS["muted"])
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: 封面 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_text(slide, 0.82, 0.95, 11.8, 0.78, "第 2 周：.data/.bss 初始化与 C/汇编混合启动", size=34, bold=True)
    add_text(slide, 0.88, 1.88, 11.5, 0.44, "单次连堂课：从段结构认知到 clear_bss 与 C/汇编混合调用", size=21, color=COLORS["blue"])
    add_text(slide, 0.9, 3.0, 11.5, 0.65, "前置：已完成第 1 周 QEMU Hello miniOS", size=24)
    add_text(slide, 0.9, 6.42, 11.6, 0.3, "实验以 tag week02-data-bss 为准，详见 docs/week02/data_bss.md。", size=13, color=COLORS["muted"])

    # ── Slide 2: 本次课安排 ──
    add_table_slide(
        prs,
        "本次课安排",
        "在第 1 周启动链路基础上，增加数据段初始化和 C/汇编混合调用。",
        None,
        ["板块", "内容"],
        [
            ("板块一", "复盘第 1 周 + 引入新问题：全局变量初始状态谁准备？"),
            ("板块二", "段结构讲解：.text / .rodata / .data / .bss"),
            ("板块三", "课堂 Demo 源码路径：main.c → start.S → linker.ld → string.S"),
            ("板块四", "编译运行与输出分析"),
            ("板块五", "学生实验安排与 AI 共学"),
            ("收尾", "核心结论、本周边界、思考题"),
        ],
        [2.4, 9.35],
        "本周不追求完整 ABI 背诵，重点是理解裸机程序必须自己准备数据段初始状态。",
    )

    # ═══════════════════════════════════════════
    # 板块一：复盘第 1 周 + 引入新问题
    # ═══════════════════════════════════════════

    add_section_slide(
        prs,
        "板块一",
        "复盘第 1 周，引入新问题",
        "学生能说出第 1 周启动链路，并能识别裸机全局变量的初始状态问题。",
    )

    add_flow_slide(
        prs,
        "第 1 周启动链路（快速回顾）",
        [
            ("_start", "boot/start.S\n设置 $sp", COLORS["blue"]),
            ("kernel_main", "进入 C 语言\n调用 printk", COLORS["green"]),
            ("printk", "逐字符输出\n调用 uart_putc", COLORS["orange"]),
            ("UART", "写 UART0_BASE\nMMIO 寄存器", COLORS["teal"]),
            ("终端", "Hello miniOS\non LoongArch64", COLORS["red"]),
        ],
        "上周解决了'CPU 如何进入 C 函数并输出字符'。",
        "板块一",
    )

    add_code_slide(
        prs,
        "引入问题：全局变量的初始状态谁准备？",
        """static char bss_buffer[16];                // 无显式初始值
static char data_message[] = "data section ok";  // 有显式初始值

void kernel_main(void)
{
    // data_message 的初始值从哪来？
    // bss_buffer[0] 是谁保证它等于 0？
}""",
        "在普通 Linux 程序里，这件事由加载器和 C 运行库完成。裸机 miniOS 中谁来负责？",
        "板块一",
    )

    add_bullet_slide(
        prs,
        "典型错误猜测 → 引出本周主题",
        [
            "❌ '硬件自动把内存清零' —— RAM 上电后是随机值。",
            "❌ '编译器会生成清零代码' —— 编译器只编译你的 C 代码，不管运行环境。",
            "❌ '全局变量天生就是 0' —— C 语义是标准规定，不是硬件行为。",
            "✅ 这些猜测恰恰说明：学生把操作系统和 C 运行库的功劳当成了'天然'。",
            "→ 本周核心：启动汇编必须在进入 C 之前，自己清零 .bss。",
        ],
        lesson="板块一",
    )

    # ═══════════════════════════════════════════
    # 板块二：段结构讲解
    # ═══════════════════════════════════════════

    add_section_slide(
        prs,
        "板块二",
        "段（section）结构讲解",
        "学生能区分 .text / .rodata / .data / .bss 四种段，理解 .bss 为什么需要清零。",
    )

    add_table_slide(
        prs,
        "程序按用途分成若干段",
        "编译链接后的镜像不是一整块无差别内存。",
        "板块二",
        ["段", "内容", "C 语言对应", "理解要求"],
        [
            (".text", "机器指令", "函数代码", "CPU 真正执行的指令"),
            (".rodata", "只读常量", "字符串字面量等", "只读，一般不修改"),
            (".data", "有初始值的全局/静态变量", 'static char s[] = "hi";', "运行时应能读到初始值"),
            (".bss", "无显式初始值的全局/静态变量", "static char buf[16];", "进入 C 前应为 0"),
        ],
        [1.2, 1.6, 3.6, 5.35],
        "本周至少要能区分 .data 和 .bss。",
    )

    add_bullet_slide(
        prs,
        "为什么要有 .bss？",
        [
            "如果把大量「全是 0 的初始值」都原样写进镜像文件，镜像会变大、浪费空间。",
            "更高效的做法：链接时只记录 .bss 需要的地址范围；运行时在进入 C 之前，把这段内存清成 0。",
            "普通 Linux 程序里这件事由操作系统加载器和 C 运行库完成。",
            "miniOS 是裸机程序，没有操作系统帮忙，所以必须由启动汇编自己做。",
        ],
        "打个比方：搬进空宿舍前统一擦一遍抽屉，而不是给每个空抽屉贴「空的」标签。",
        lesson="板块二",
    )

    add_bullet_slide(
        prs,
        "两个必须强调的易混淆点",
        [
            "① 'C 语言说初始为 0' ≠ '硬件自动变成 0'。C 语义是标准规定，但 CPU 和内存不会自动实现。裸机中必须有人真正执行「写 0」这个动作。",
            '② __bss_start / __bss_end 不是 C 全局数组。它们是链接脚本提供的地址符号——就像家具摆放图上的尺寸线，告诉工人"从哪清到哪"，而不是实物家具。',
            "学生容易用 C 变量的思维理解链接脚本符号——以为 __bss_start 是一个存着地址值的 int 变量。要强调它是在链接阶段由链接器计算出的地址标签。",
        ],
        lesson="板块二",
    )

    # ═══════════════════════════════════════════
    # 板块三：课堂 Demo 源码路径
    # ═══════════════════════════════════════════

    add_section_slide(
        prs,
        "板块三",
        "课堂 Demo：源码路径",
        "学生能沿着四个关键文件，理解 .data/.bss 验证和 clear_bss 的完整执行路径。",
    )

    add_bullet_slide(
        prs,
        "本周重点文件（按阅读顺序）",
        [
            "① kernel/main.c —— 先看「验证了什么」：data_message 与 bss_buffer。",
            "② boot/start.S —— 再看「进入 C 前做了什么」：clear_bss 实现。",
            "③ kernel/linker.ld —— 再看「清零边界从哪来」：__bss_start / __bss_end。",
            "④ lib/string.S + include/string.h —— 最后看「C 如何调用汇编库」。",
        ],
        lesson="板块三",
    )

    add_code_annotated_slide(
        prs,
        "kernel/main.c —— 验证逻辑",
        """static char bss_buffer[16];
static char data_message[] = "data section ok";

void kernel_main(void) {
    char buf[32];
    printk("Hello miniOS on LoongArch64\\n");  // ① 第 1 周链路

    memset(buf, 0, sizeof(buf));               // ② 清零局部 buf
    memcpy(buf, data_message, strlen(data_message)); // ③ 从 .data 复制
    printk(buf); printk("\\n");                 // ④ 输出 data section ok

    if (bss_buffer[0] == 0)                    // ⑤ 检查 .bss 是否已清零
        printk("bss section cleared\\n");

    printk("week1-week2 check done\\n");        // ⑥ 阶段标记
    while (1) { __asm__ volatile("idle 0"); }   // ⑦ 停机
}""",
        [
            "① 保留第 1 周 Hello，验证启动链路仍有效。",
            "②③④ 验证 .data 可读：用汇编实现的 memset/memcpy/strlen 把 data_message 复制到栈上再输出。",
            "⑤ 验证 .bss 已清零：只有 clear_bss 在 kernel_main() 之前执行，这里才一定成立。",
            "⑥⑦ 阶段完成标记 + 停机循环，防止 CPU 跑飞。",
            "这里用 memcpy 而非直接 printk(data_message)，是为了练习完整的 C/汇编库调用链路。",
        ],
        "先看 main.c 知道「验证了什么」，再去看启动汇编「怎么做到的」。",
        "板块三",
    )

    add_code_annotated_slide(
        prs,
        "boot/start.S —— _start 完整职责",
        """    .section .text.boot, "ax"
    .globl _start

_start:
    la.global   $sp, boot_stack_top     @ ① 设置内核栈
    bl          clear_bss               @ ② 清零 .bss ← 本周新增
    bl          kernel_main             @ ③ 进入 C 语言内核

halt:
    idle        0
    b           halt""",
        [
            "① la.global $sp, boot_stack_top：CPU 上电后 $sp 是垃圾值，必须先指向预留栈顶。",
            "② bl clear_bss：本周新增关键步骤——清零 .bss 段，保证未初始化全局变量从 0 开始。",
            "③ bl kernel_main：调用 C 函数。注意 clear_bss 必须在 kernel_main 之前——顺序不可颠倒。",
            "halt 循环：防止 kernel_main 意外返回后 CPU 从未知内存取指执行。",
        ],
        "对比第 1 周：在设置 $sp 和 kernel_main 之间多了 bl clear_bss。",
        "板块三",
    )

    add_code_annotated_slide(
        prs,
        "clear_bss 逐行讲解",
        """clear_bss:
    la.global   $t0, __bss_start        @ t0 = 清零起点
    la.global   $t1, __bss_end          @ t1 = 清零终点

1:
    beq         $t0, $t1, 2f            @ t0 == t1? 跳出循环
    st.b        $zero, $t0, 0           @ *t0 = 0（逐字节写 0）
    addi.d      $t0, $t0, 1             @ t0++
    b           1b                       @ 继续循环
2:
    jr          $ra                      @ 返回""",
        [
            "自然语言翻译：t0 = __bss_start; t1 = __bss_end; while (t0 != t1) { *t0 = 0; t0++; } return;",
            "__bss_start / __bss_end 来自链接脚本，不是在这里定义的。",
            "清零区间 [__bss_start, __bss_end)，左闭右开。",
            "$zero 恒为 0，st.b 逐字节存储——教学版追求清晰易懂，不追求性能。",
            "后续周次再讨论按 8 字节批量清零的优化方案。",
        ],
        "这是 miniOS 真实的 clear_bss 实现，不是教学示意代码。",
        "板块三",
    )

    add_code_annotated_slide(
        prs,
        "kernel/linker.ld —— 清零边界从哪来",
        """ENTRY(_start)
SECTIONS {
    . = 0x9000000000200000;

    .text : ALIGN(4K) { KEEP(*(.text.boot)) *(.text .text.*) }
    .rodata : ALIGN(4K) { *(.rodata .rodata.*) }
    .data : ALIGN(4K) { *(.data .data.*) }

    __bss_start = .;                     @ ← 在 .bss 前记录当前地址
    .bss : ALIGN(4K) {
        *(.bss .bss.*)
        *(COMMON)
    }
    __bss_end = .;                       @ ← 在 .bss 后记录当前地址
}""",
        [
            "'.' 是链接器的位置计数器，表示当前地址。",
            "__bss_start = . 在 .bss 段开始前记录当前地址；__bss_end = . 在段结束后记录。",
            "这两个符号是链接时生成的地址标签，不是 C 变量。汇编中用 la.global 加载的是地址值本身。",
            "可以比作尺子上的两个刻度——它们是位置标记，不占用 .bss 区间内的额外空间。",
            "ENTRY(_start) 告诉链接器入口符号是 _start，不是 main。",
        ],
        "链接脚本决定了代码和数据的布局，__bss_start/__bss_end 是其中的关键标签。",
        "板块三",
    )

    add_table_slide(
        prs,
        "C 与汇编混合调用",
        "本周只需建立三个认知：同名符号、统一约定、一起链接。",
        "板块三",
        ["环节", "文件", "关键点"],
        [
            ("声明", "include/string.h", "memset / memcpy / strlen 的函数原型，供 C 编译器生成调用代码"),
            ("定义", "lib/string.S", ".globl 导出同名全局符号，按 ABI 约定在寄存器中收发参数和返回值"),
            ("链接", "Makefile", "lib/string.S 编进 SRCS_S，与其他 .o 一起链入 minios.elf"),
        ],
        [1.5, 2.8, 7.45],
        "本周不要求背完整 LoongArch ABI，只需知道 '有约定' 即可。",
    )

    add_code_slide(
        prs,
        "lib/string.S —— memset 为例",
        """    .globl memset
memset:
    move        $t0, $a0          @ 保存原始 dst（返回值）
    beqz        $a2, 2f           @ n == 0 直接返回

1:
    st.b        $a1, $a0, 0       @ *dst = value
    addi.d      $a0, $a0, 1       @ dst++
    addi.d      $a2, $a2, -1      @ n--
    bnez        $a2, 1b           @ n != 0 继续

2:
    move        $a0, $t0           @ 返回值 = 原始 dst
    jr          $ra""",
        '参数：$a0=dst, $a1=value, $a2=n；返回值放在 $a0。不要求背完整 ABI，建立"有约定"的意识即可。',
        "板块三",
    )

    # ═══════════════════════════════════════════
    # 板块四：编译运行与输出分析
    # ═══════════════════════════════════════════

    add_section_slide(
        prs,
        "板块四",
        "编译、运行与输出分析",
        "学生能独立完成编译运行，并逐行解释四行验收输出的含义。",
    )

    add_code_slide(
        prs,
        "课堂演示命令",
        """# 环境检查
sh scripts/check-env.sh

# 编译
make clean
make

# 运行
make run

# 退出 QEMU：先按 Ctrl-a，松开后再按 x""",
        "预期输出四行，见下一页。未执行过 make run 就不能写「已通过」，工具链缺失要记录「未执行」。",
        "板块四",
    )

    add_table_slide(
        prs,
        "验收输出逐行分析",
        None,
        "板块四",
        ["输出行", "验证了什么", "如果不出现说明什么"],
        [
            ("Hello miniOS on LoongArch64", "第 1 周启动链路仍有效", "启动链路有问题，先回到第 1 周排查"),
            ("data section ok", ".data 段可读，C 成功调用汇编库", "data_message 可能不在 .data，或 memcpy/strlen 有误"),
            ("bss section cleared", "clear_bss 已执行，.bss 段为 0", "clear_bss 未调用、清零区间错误、或 bss_buffer 未落在 .bss"),
            ("week1-week2 check done", "检查路径执行到末尾", "前面某一步提前失败或卡住"),
        ],
        [3.3, 4.0, 4.45],
        "若某行缺失，按 main.c → start.S → linker.ld 顺序逐文件排查，不要盲目重编。",
    )

    add_flow_slide(
        prs,
        "第 2 周完整链路",
        [
            ("_start", "设置 $sp\nla.global", COLORS["blue"]),
            ("clear_bss", "清零 .bss\n[__bss_start,\n__bss_end)", COLORS["orange"]),
            ("kernel_main", "验证 .data\n验证 .bss\n输出四行", COLORS["green"]),
            ("memset/\nmemcpy/\nstrlen", "汇编库函数\nC 侧声明调用", COLORS["teal"]),
            ("UART", "终端显示\n四行输出", COLORS["red"]),
        ],
        "第 1 周搭栈，第 2 周搭数据初始状态——每一周都在为 C 代码添加基础设施。",
        "板块四",
    )

    # ═══════════════════════════════════════════
    # 板块五：学生实验安排与 AI 共学
    # ═══════════════════════════════════════════

    add_section_slide(
        prs,
        "板块五",
        "学生实验安排与 AI 共学",
        "学生能独立创建实验分支、完成环境检查和编译运行，并知晓 AI 使用边界。",
    )

    add_code_slide(
        prs,
        "学生实验任务",
        """# 1. 建立实验分支
git fetch --tags
git switch -c my-week02-lab week02-data-bss
git branch --show-current
git describe --tags --always

# 2. 环境检查
sh scripts/check-env.sh           # 如实记录，不要写"已通过"

# 3. 编译运行
make clean && make && make run     # 记录四行输出是否齐全

# 4. 源码阅读（必做）
#    - data_message vs bss_buffer 区别
#    - clear_bss 执行流程（文本流程图）
#    - __bss_start / __bss_end 来自哪里
#    - C 为何能调用 lib/string.S

# 5. 可选加深：objdump -d 观察反汇编、nm 查看符号地址""",
        "实验报告模板见 docs/week02/data_bss.md 第 17 节。",
        "板块五",
        font_size=11,
    )

    add_bullet_slide(
        prs,
        "AI 共学：允许 vs 不允许",
        [
            "✅ 请 AI 解释 clear_bss 的执行流程。",
            "✅ 请 AI 画出 _start → clear_bss → kernel_main() 的文本流程图。",
            "✅ 把真实的 make / make run 错误信息发给 AI 分析原因。",
            "✅ 让 AI 辅助整理实验报告的结构和文字表达。",
            "❌ 让 AI 编造「已跑通」的串口输出。",
            "❌ 不读源码，直接让 AI 代写全部分析后原样提交。",
            "❌ 把 AI 猜测结果写成实测结果。",
        ],
        "第 2 周起 AI 可以做更多，但角色仍是'助教'而非'替身'。实验报告输出必须来自真实终端。",
        lesson="板块五",
    )

    # ═══════════════════════════════════════════
    # 收尾
    # ═══════════════════════════════════════════

    add_section_slide(
        prs,
        "收尾",
        "核心结论、本周边界与思考题",
        "学生能用自己的话总结本周核心结论，明确哪些内容本周不讲、留给后续周次。",
    )

    add_bullet_slide(
        prs,
        "本周核心结论",
        [
            "裸机 C 程序不是「天然能正确运行」——启动汇编必须准备最小运行环境。",
            "三条准备工作：① 设置栈  ② 清零 .bss  ③ 再进入 kernel_main。",
            ".data 有初值（镜像保存），.bss 无显式初值但语义为 0（启动时清零）。",
            "C 与汇编通过统一符号名和链接过程协同工作：.h 声明 + .S 定义 + Makefile 一起链。",
            "第 1 周搭了'栈'，第 2 周搭了'数据初始状态'——后面每一周都会在这个舞台上添加新基础设施。",
        ],
        lesson="收尾",
    )

    add_bullet_slide(
        prs,
        "本周边界",
        [
            "不要求从零实现高性能 memset/memcpy（第 10 周深入）。",
            "不要求背诵完整 LoongArch ABI 寄存器保存规则。",
            "不涉及开发板移植——仍然是 QEMU virt。",
            "kernel/exception.c、kernel/syscall.c 属于后续周次铺垫，本周不展开。",
            "lib/string.S 中 memset 使用逐字节循环——教学上先保证清晰，性能优化后续再讲。",
        ],
        lesson="收尾",
    )

    add_bullet_slide(
        prs,
        "思考题",
        [
            "① .data 和 .bss 的区别是什么？结合 data_message 与 bss_buffer 说明。",
            "② 为什么 .bss 清零必须在 kernel_main() 之前？放在之后会怎样？",
            "③ __bss_start 和 __bss_end 是谁提供的？它们是 C 变量吗？",
            "④ C 代码为什么能调用 lib/string.S 中的汇编函数？从声明、定义、链接三点回答。",
            "⑤ 如果只有 Hello 没有 .data/.bss 检查结果，应按什么顺序排查哪些文件？",
            "⑥ 普通 Linux 程序由谁负责准备 .data 和清零 .bss？和 miniOS 有什么不同？",
        ],
        "作业：完成实验并撰写实验报告（模板见 data_bss.md §17），输出必须来自真实命令执行。",
        lesson="收尾",
    )

    # ── 最后：第 2 周 vs 第 1 周对比 ──
    add_flow_slide(
        prs,
        "第 1 周 → 第 2 周：路径对比",
        [
            ("第 1 周", "_start\n→ 设置 sp\n→ kernel_main\n→ Hello", COLORS["blue"]),
            ("第 2 周新增", "clear_bss\n.data 验证\n.bss 验证\nC/汇编混合", COLORS["orange"]),
            ("第 3 周预告", "分支与循环\n字符串输出\n控制流组织", COLORS["muted"]),
        ],
        "每一周在上一周的基础上增加一层基础设施。第 3 周将在本节基础上引入控制流。",
        "收尾",
    )

    prs.save(OUT)
    print(OUT)
    print(f"slides={len(prs.slides)}")


if __name__ == "__main__":
    build()
