from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "week01" / "week01_qemu_hello_course.pptx"

COLORS = {
    "ink": RGBColor(34, 37, 43),
    "muted": RGBColor(91, 99, 112),
    "bg": RGBColor(248, 249, 251),
    "line": RGBColor(216, 221, 230),
    "blue": RGBColor(28, 83, 160),
    "green": RGBColor(42, 126, 78),
    "orange": RGBColor(184, 95, 34),
    "teal": RGBColor(0, 116, 130),
    "red": RGBColor(164, 68, 64),
    "code_bg": RGBColor(32, 36, 44),
    "code_fg": RGBColor(238, 241, 246),
    "white": RGBColor(255, 255, 255),
}


def apply_font(run, size=22, bold=False, color=None, name="Microsoft YaHei"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]


def set_speaker_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text
    return slide


def set_background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["bg"]


def add_text(slide, x, y, w, h, body, size=22, bold=False, color=None, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = body
    apply_font(r, size=size, bold=bold, color=color, name=font_name)
    return box


def add_title(slide, heading, subtitle=None, lesson=None):
    if lesson:
        add_text(slide, 0.7, 0.28, 2.4, 0.3, lesson, size=13, bold=True, color=COLORS["blue"])
    add_text(slide, 0.7, 0.55, 12.0, 0.58, heading, size=30, bold=True)
    if subtitle:
        add_text(slide, 0.72, 1.14, 11.9, 0.35, subtitle, size=14, color=COLORS["muted"])
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(1.48), Inches(11.95), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.fill.background()


def add_bullet_slide(prs, heading, items, subtitle=None, lesson=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, heading, subtitle, lesson)
    box = slide.shapes.add_textbox(Inches(0.95), Inches(1.85), Inches(11.45), Inches(5.15))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(9)
        r = p.add_run()
        r.text = item
        apply_font(r, size=20 if len(item) < 34 else 17)
    return slide


def add_code_slide(prs, heading, code_text, subtitle=None, lesson=None, font_size=13):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, heading, subtitle, lesson)
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.82),
        Inches(1.75),
        Inches(11.72),
        Inches(5.08),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["code_bg"]
    shape.line.color.rgb = COLORS["code_bg"]
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.16)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = code_text
    apply_font(r, size=font_size, color=COLORS["code_fg"], name="Consolas")
    return slide


def add_code_annotated_slide(prs, heading, code_text, annotations, subtitle=None, lesson=None, font_size=12):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, heading, subtitle, lesson)
    code_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(1.75),
        Inches(6.6),
        Inches(5.08),
    )
    code_shape.fill.solid()
    code_shape.fill.fore_color.rgb = COLORS["code_bg"]
    code_shape.line.color.rgb = COLORS["code_bg"]
    tf = code_shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.14)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = code_text
    apply_font(r, size=font_size, color=COLORS["code_fg"], name="Consolas")

    box = slide.shapes.add_textbox(Inches(7.5), Inches(1.85), Inches(5.15), Inches(4.95))
    note_tf = box.text_frame
    note_tf.clear()
    note_tf.word_wrap = True
    for i, item in enumerate(annotations):
        para = note_tf.paragraphs[0] if i == 0 else note_tf.add_paragraph()
        para.space_after = Pt(8)
        run = para.add_run()
        run.text = item
        apply_font(run, size=14)
    return slide


def add_section_slide(prs, lesson, heading, goal):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_text(slide, 0.85, 1.1, 2.3, 0.35, lesson, size=15, bold=True, color=COLORS["blue"])
    add_text(slide, 0.85, 1.72, 11.4, 0.78, heading, size=36, bold=True)
    add_text(slide, 0.9, 2.86, 10.9, 0.95, goal, size=23, color=COLORS["muted"])
    add_text(slide, 0.9, 6.55, 11.4, 0.28, "本节结束时，学生要能用自己的话解释核心概念。", size=13, color=COLORS["muted"])
    return slide


def add_flow_slide(prs, heading, steps, subtitle=None, lesson=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, heading, subtitle, lesson)
    x = 0.62
    width = 2.18 if len(steps) == 5 else 2.45
    gap = 0.25
    for i, (name, body, color) in enumerate(steps):
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.35), Inches(width), Inches(1.75))
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLORS["white"]
        rect.line.color.rgb = COLORS["line"]
        add_text(slide, x + 0.16, 2.53, width - 0.32, 0.32, name, size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.18, 2.96, width - 0.36, 0.82, body, size=13, align=PP_ALIGN.CENTER)
        x += width + gap
        if i < len(steps) - 1:
            add_text(slide, x - 0.05, 3.0, 0.25, 0.3, "→", size=20, bold=True, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    return slide


def add_table_slide(prs, heading, subtitle, lesson, headers, rows, col_widths, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, heading, subtitle, lesson)
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.8), Inches(1.9), Inches(11.75), Inches(4.25))
    table = table_shape.table
    for idx, width in enumerate(col_widths):
        table.columns[idx].width = Inches(width)
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["blue"]
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                apply_font(r, size=13, bold=True, color=COLORS["white"])
    for row_idx, row in enumerate(rows, start=1):
        for col, value in enumerate(row):
            cell = table.cell(row_idx, col)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["white"]
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    apply_font(r, size=12)
    add_text(slide, 0.82, 6.55, 11.7, 0.35, footer, size=13, color=COLORS["muted"])
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_text(slide, 0.82, 0.95, 11.8, 0.78, "第 1 周：从 0 启动一个 LoongArch miniOS", size=34, bold=True)
    add_text(slide, 0.88, 1.88, 11.5, 0.44, "单次连堂课：从为什么学汇编到 QEMU Hello miniOS", size=21, color=COLORS["blue"])
    add_text(slide, 0.9, 3.0, 11.5, 0.65, "对象：已系统学习 C 语言的大二计算机类学生", size=24)
    add_text(slide, 0.9, 6.42, 11.6, 0.3, "实验以本仓库 miniOS 代码与 docs/week01_qemu_hello.md 为准。", size=13, color=COLORS["muted"])

    add_table_slide(
        prs,
        "本次课安排",
        "四节课内容压缩进一次连堂课，按板块推进。",
        None,
        ["板块", "内容"],
        [
            ("板块一", "为什么学习汇编：C 语言经验如何连接到 CPU 执行"),
            ("板块二", "LoongArch 最小知识包：寄存器、指令"),
            ("板块三", "从 C 到可执行文件：编译链接、Makefile 语法、操作系统做什么、裸机差异"),
            ("板块四", "QEMU Hello miniOS 实验：boot.S 详解与实际运行"),
            ("收尾", "本周边界、思考题与作业"),
        ],
        [2.4, 9.35],
        "第一周不追求掌握全部指令，而是建立正确的底层执行模型。",
    )

    add_flow_slide(
        prs,
        "从 C 代码到 CPU 执行",
        [
            ("C 语言", "变量\n函数\n循环", COLORS["green"]),
            ("汇编语言", "寄存器\n跳转\n访存", COLORS["blue"]),
            ("机器指令", "二进制编码\nCPU 取指执行", COLORS["orange"]),
            ("硬件行为", "运算\n读写内存\n访问外设", COLORS["teal"]),
        ],
        "汇编课的目标不是背指令，而是看清 C 程序如何落到机器执行。",
        "板块一",
    )
    add_bullet_slide(
        prs,
        "打个比方：CPU 是只懂机器语言的专家",
        [
            "CPU 就像一位只精通『机器语言』的专家，完全听不懂我们写的 C 语言。",
            "编译器/汇编器就是翻译：先把 C 语言译成汇编语言——这是一份贴近专家表达习惯的逐句译稿。",
            "汇编语言再被汇编器翻译成机器指令，这才是专家真正执行的语言，只剩 0 和 1。",
            "我们学汇编，就是学会读懂这份『翻译中间稿』，看清 C 程序最终变成了什么。",
        ],
        lesson="板块一",
    )
    add_bullet_slide(
        prs,
        "用 C 语言经验理解汇编",
        [
            "C 里的变量：很多时候会临时放在寄存器或内存里。",
            "C 里的 if / while：底层会变成比较和跳转指令。",
            "C 里的函数调用：底层涉及参数寄存器、返回地址和栈。",
            "C 里的指针：底层就是地址，访存指令按地址读写数据。",
        ],
        lesson="板块一",
    )
    add_code_annotated_slide(
        prs,
        "真实例子逐行拆解：add() 函数",
        """int add(int a, int b)
{
    return a + b;
}

# 对应 LoongArch 汇编（-O0，教学示意）
add:
    addi.d  $sp, $sp, -16      # 开辟栈空间
    st.w    $a0, $sp, 12       # 参数 a 存入栈
    st.w    $a1, $sp, 8        # 参数 b 存入栈
    ld.w    $t0, $sp, 12       # 取回 a
    ld.w    $t1, $sp, 8        # 取回 b
    add.w   $a0, $t0, $t1      # 计算 a + b，结果放入返回值寄存器
    addi.d  $sp, $sp, 16       # 释放栈空间
    jirl    $zero, $ra, 0      # 返回调用者""",
        [
            "① $a0、$a1 是参数寄存器，对应 a、b——翻译稿里已经用上了『专家的方言』。",
            "② 出现栈上存取是因为 -O0 未优化；优化后可能直接用寄存器完成，不落栈。",
            "③ add.w 是唯一『做加法』的指令，其余都是数据搬运。",
            "④ 返回值放在 $a0。",
            "⑤ jirl $zero, $ra, 0 配合 $ra 完成函数返回。",
            "教学示意汇编，非本机实测编译输出：本机未装 loongarch64-linux-gnu-gcc。",
        ],
        "这就是刚才『翻译』比喻里的翻译稿：C 函数被翻译成了这份汇编。",
        "板块一",
    )
    add_bullet_slide(
        prs,
        "课堂活动：徒手推演 max()",
        [
            "任务：int max(int a, int b) { return a > b ? a : b; }",
            "独立思考：不写真实汇编语法，只用文字描述 CPU 需要做哪几步。",
            "同桌讨论：重点讨论『比较』和『选择』在硬件层面可能对应什么动作。",
            "抽 1-2 人口头说明，为板块二的分支跳转指令做铺垫。",
        ],
        lesson="板块一",
    )

    add_table_slide(
        prs,
        "LoongArch 通用寄存器最小表",
        "第一周先认编号和常用 ABI 别名。",
        "板块二",
        ["编号", "ABI 别名", "主要用途", "课堂说明"],
        [
            ("r0", "zero", "常量 0", "硬件恒为 0"),
            ("r1", "ra", "返回地址", "bl 调用函数时写入"),
            ("r3", "sp", "栈指针", "进入 C 前必须设置"),
            ("r4-r11", "a0-a7", "参数寄存器", "函数参数，前两个也可作返回值"),
            ("r12-r20", "t0-t8", "临时寄存器", "短期计算常用"),
            ("r22-r31", "fp/s0-s9", "保存寄存器", "函数调用约定后续展开"),
        ],
        [1.35, 2.1, 2.3, 5.95],
        "讲解顺序：先用 r 编号建立硬件视角，再用 ABI 别名阅读源码。",
    )
    add_bullet_slide(
        prs,
        "打个比方：寄存器就像口袋",
        [
            "寄存器就像你随身衣服上的口袋：数量有限，但掏取速度最快。",
            "内存就像一个大仓库：能装的东西多得多，但要跑一趟去取，明显更慢。",
            "LoongArch 只有 32 个『口袋』（r0-r31），装不下的数据要临时存进『仓库』（内存/栈）。",
            "这就是为什么 boot/start.S 一开始要先准备好『仓库入口』（设置 $sp 指向栈）——口袋不够用时才用得上它。",
        ],
        lesson="板块二",
    )
    add_code_slide(
        prs,
        "基础指令例子",
        """addi.d  r3, r3, -32       # sp = sp - 32
st.d    r1, r3, 24        # 保存 ra
ld.d    r1, r3, 24        # 恢复 ra
bne     r12, r0, L        # r12 != 0 时跳转
bl      func              # 调用函数，返回地址写入 r1/ra
jirl    r0, r1, 0         # 返回到 ra 指向的位置""",
        "先用 r 编号讲硬件动作，再映射到 ABI 别名；GNU 汇编器同时接受 $sp/$ra 这类别名，$sp 就是 r3。",
        "板块二",
    )

    add_flow_slide(
        prs,
        "编译流程，用到本实验里",
        [
            (".c/.S", "源代码\nC 和汇编", COLORS["green"]),
            ("编译", "生成目标文件\n处理 include", COLORS["blue"]),
            ("链接", "合成 ELF\n确定地址入口", COLORS["orange"]),
            ("QEMU", "加载内核\n开始执行", COLORS["teal"]),
        ],
        "第 1 周只需要说出每一步的作用。",
        "板块三",
    )
    add_code_slide(
        prs,
        "Makefile 语法速览",
        """# Makefile 最小语法结构（通用示例，不是本项目代码）
目标: 依赖1 依赖2
\t命令              # 命令前必须是 Tab，不能是空格

# 例子：把 hello.c 编译成 hello.o
hello.o: hello.c
\tgcc -c hello.c -o hello.o

CC := gcc            # := 立即展开赋值
CFLAGS ?= -Wall       # ?= 只在变量还没被设置时才赋值

.PHONY: clean         # 声明伪目标，clean 不是一个真实文件
clean:
\trm -f *.o""",
        "先认识目标/依赖/命令三段式，再看本项目 Makefile 是怎么用这些语法的。",
        "板块三",
        font_size=12,
    )
    add_code_slide(
        prs,
        "Makefile 的核心含义",
        """CROSS_COMPILE ?= loongarch64-linux-gnu-
CC      := $(CROSS_COMPILE)gcc
OBJCOPY := $(CROSS_COMPILE)objcopy
QEMU    ?= qemu-system-loongarch64

TARGET  := build/minios.elf
LDFLAGS := -T kernel/linker.ld -nostdlib -static""",
        "不能用宿主机 x86_64 gcc 冒充 LoongArch 工具链。",
        "板块三",
    )
    add_code_slide(
        prs,
        "链接脚本先讲三件事",
        """ENTRY(_start)

SECTIONS
{
    . = 0x9000000000200000;

    .text : { *(.text.boot) *(.text*) }
    .rodata : { *(.rodata*) }
    .data : { *(.data*) }
    .bss : { *(.bss*) *(COMMON) }
}""",
        "入口是谁、从哪个地址开始、各类内容放在哪里。",
        "板块三",
    )
    add_bullet_slide(
        prs,
        "操作系统平时帮你做了什么",
        [
            "帮你把程序加载到内存、分配好栈和堆，你的 C 代码打开就能跑。",
            "帮你管理多个程序同时运行（进程调度），你不用关心 CPU 什么时候轮到你。",
            "帮你把 printf、malloc 这些函数背后的系统调用接好，你才能直接调用。",
            "帮你管理磁盘文件、网络、显示器这些硬件，你只需要调用统一的接口。",
        ],
        "普通 C 程序能“写完就跑”，全靠操作系统在背后先把环境搭好。",
        lesson="板块三",
    )
    add_bullet_slide(
        prs,
        "普通程序和裸机程序的启动差异",
        [
            "普通 Linux 程序：操作系统和 C 运行库准备运行环境。",
            "miniOS 裸机程序：没有 libc，没有操作系统帮忙调用 main。",
            "所以我们需要 _start、栈、链接脚本和串口输出。",
            "这正是第 1 周 Hello miniOS 的教学价值。",
        ],
        lesson="板块三",
    )

    add_bullet_slide(
        prs,
        "第 1 周核心文件",
        [
            "boot/start.S：启动入口，设置栈，跳转 kernel_main。",
            "kernel/linker.ld：规定入口地址和段布局。",
            "kernel/main.c：C 语言内核入口，调用 printk。",
            "kernel/printk.c：最小字符串输出。",
            "include/printk.h、include/uart.h：函数原型和 UART 地址。",
            "Makefile：编译、运行和清理命令。",
        ],
        lesson="板块四",
    )
    add_code_annotated_slide(
        prs,
        "boot/start.S 主路径逐行讲解",
        """    .section .text.boot, "ax"
    .globl _start

_start:
    # 设置内核栈，$sp 是 r3 的 ABI 别名
    la.global   $sp, boot_stack_top

    # 进入 C 语言内核主函数
    bl          kernel_main

halt:
    idle        0
    b           halt""",
        [
            "① .section .text.boot 把这段代码放到链接脚本指定的启动区域，保证它在最前面。",
            "② .globl _start 让链接器能找到这个入口符号，对应链接脚本里的 ENTRY(_start)。",
            "③ la.global $sp, boot_stack_top：CPU 上电后 $sp 是垃圾值，必须先指向预留的栈顶。",
            "④ bl kernel_main：调用 C 函数，同时把返回地址写入 $ra——但 kernel_main 不会返回。",
            "⑤ halt 循环：防止 kernel_main 意外返回后，CPU 从未知内存继续取指执行。",
        ],
        "这是 miniOS 里真实的 boot/start.S 主路径，不是教学示意代码。",
        "板块四",
    )
    add_code_slide(
        prs,
        "建议第 1 周使用的最小 kernel_main",
        """#include "printk.h"

void kernel_main(void)
{
    printk("Hello miniOS on LoongArch64\\n");

    while (1) {
        __asm__ volatile("idle 0");
    }
}""",
        "当前 master 有第 2 周检查代码；课堂先抽出这条最小路径。",
        "板块四",
    )
    add_code_slide(
        prs,
        "printk 到 UART",
        """void uart_putc(char ch)
{
    volatile unsigned char *uart =
        (volatile unsigned char *)UART0_BASE;

    *uart = (unsigned char)ch;
}

void printk(const char *s)
{
    while (*s) {
        uart_putc(*s++);
    }
}""",
        "UART0_BASE 是 QEMU virt 平台地址，不代表 2K0300 开发板地址。",
        "板块四",
    )
    add_bullet_slide(
        prs,
        "打个比方：一场接力赛",
        [
            "QEMU 是裁判员：把 minios.elf 放上起跑线（加载进内存），鸣枪示意从入口地址开始（CPU 取指）。",
            "_start 是第一棒选手：上场先系好鞋带（设置 $sp 指向栈），再把接力棒交给下一棒（bl kernel_main）。",
            "kernel_main 是第二棒主力选手：真正跑出比赛内容（调用 printk）。",
            "printk/uart_putc 是选手把成绩喊给场边记录员（把字节写进 UART0_BASE 寄存器）。",
            "QEMU 终端是记分牌：把喊出的成绩显示给所有观众看（终端里出现 Hello miniOS）。",
            "halt 循环是比赛结束后运动员原地休息——裸机没有『下一场』可跑，不会返回。",
        ],
        "对照下一页的流程图和文字版逐步说明，这里先建立直观印象。",
        lesson="板块四",
    )
    add_flow_slide(
        prs,
        "第 1 周最小运行链路",
        [
            ("make", "交叉编译\n生成 ELF", COLORS["blue"]),
            ("QEMU", "模拟 virt\n加载内核", COLORS["teal"]),
            ("_start", "设置 sp\n调用 C", COLORS["orange"]),
            ("kernel_main", "调用 printk\n输出字符串", COLORS["green"]),
            ("UART", "终端显示\nHello", COLORS["red"]),
        ],
        "这条链路跑通，才进入第 2 周。",
        "板块四",
    )
    add_bullet_slide(
        prs,
        "裸机启动逐步说明",
        [
            "① QEMU 把 build/minios.elf 加载进虚拟内存，按 ELF 头找到入口地址。",
            "② CPU 从入口地址（_start）取出第一条指令开始执行，此时还没有任何 C 环境。",
            "③ _start 把 $sp 指向预留的栈空间——没有这一步，C 函数完全不能用。",
            "④ bl kernel_main 跳转进 C 代码，从这里开始才是我们熟悉的 C 语言世界。",
            "⑤ kernel_main 调用 printk，printk 逐字符调用 uart_putc。",
            "⑥ uart_putc 直接向 UART0_BASE 这个内存地址写字节——这不是普通内存，是外设寄存器。",
            "⑦ QEMU 把这次写操作翻译成终端输出，我们才看到 Hello miniOS。",
            "⑧ kernel_main 结束后代码进入 halt 死循环，因为裸机没有“返回到操作系统”这回事。",
        ],
        "对照上一页的流程图，这里是文字版逐步说明。",
        lesson="板块四",
    )
    add_code_slide(
        prs,
        "实验命令",
        """# 环境检查
sh scripts/check-env.sh

# 编译
make clean
make

# 运行
make run""",
        "预期串口输出：Hello miniOS on LoongArch64。未执行过 make run 就不能写“已通过”，工具链缺失要记录“未执行”和失败原因。",
        "板块四",
    )

    add_bullet_slide(
        prs,
        "本周边界",
        [
            "不讲完整操作系统，不讲进程、文件系统、虚拟内存。",
            "不直接适配 2K0300 开发板。",
            "不把未实测结果写成已通过——工具链缺失就记录“未执行”。",
            "当前 master 已包含第 2 周 .data/.bss 检查代码和 clear_bss，第 1 周课堂只讲 Hello 链路，这部分下周展开。",
            "lib/string.S、kernel/exception.c、kernel/syscall.c 属于后续周次铺垫，本周不展开。",
            "AI 共学：可以让 AI 解释代码、画流程图，不允许让 AI 直接生成实验代码或代替写实验报告。",
        ],
        lesson="收尾",
    )
    add_bullet_slide(
        prs,
        "思考题与作业",
        [
            "思考：CPU 第一条指令在哪里？为什么裸机程序不是从 main() 开始？",
            "思考：为什么进入 C 函数前要设置 sp？为什么 miniOS 不能直接使用 printf？",
            "思考：为什么先 QEMU 再开发板？",
            "作业：画出 miniOS 第 1 周执行路径图，把 Hello 字符串改成自己的姓名和学号重新编译运行。",
            "作业：解释 boot/start.S 中 sp、bl、b、idle 的作用，提交环境检查结果和真实运行截图或日志。",
        ],
        lesson="收尾",
    )

    prs.save(OUT)
    print(OUT)
    print(f"slides={len(prs.slides)}")


if __name__ == "__main__":
    build()
