"""Generate a standalone 2-slide pptx with a Makefile syntax primer.

Meant to be manually copy-pasted into week01_qemu_hello_course.pptx
between slide 17 ("编译流程，一路走到 QEMU 能跑起来") and slide 18
("Makefile 的核心含义：别找错翻译官"), which currently jumps straight
into project-specific Makefile variables without ever explaining the
generic target/dependency/command syntax.

Clones existing slide styles (icon-row list = slide 6, command panel =
slide 18) from the current deck via raw XML copy so the new slides
match visually, then deletes every other slide from the output copy.
"""

import copy
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "docs" / "week01" / "week01_qemu_hello_course.pptx"
OUT = ROOT / "docs" / "week01" / "week01_makefile_syntax_insert.pptx"


def duplicate_slide(prs, index):
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)

    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    for shp in source.shapes:
        new_el = copy.deepcopy(shp._element)
        dest.shapes._spTree.append(new_el)

    for shp in dest.shapes:
        blip = shp._element.find(".//" + qn("a:blip"))
        if blip is not None:
            old_rid = blip.get(qn("r:embed"))
            if old_rid:
                image_part = source.part.related_part(old_rid)
                new_rid = dest.part.relate_to(
                    image_part,
                    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
                )
                blip.set(qn("r:embed"), new_rid)

    return dest


def set_texts(slide, text_by_shape_name):
    for shp in slide.shapes:
        if shp.name in text_by_shape_name and shp.has_text_frame:
            new_text = text_by_shape_name[shp.name]
            tf = shp.text_frame
            lines = new_text.split("\n")
            base_p = tf.paragraphs[0]
            if not base_p.runs:
                continue
            base_run = base_p.runs[0]
            base_run.text = lines[0]
            for extra_run in list(base_p.runs[1:]):
                extra_run.text = ""
            for p in list(tf.paragraphs[1:]):
                p._p.getparent().remove(p._p)
            for line in lines[1:]:
                new_p = copy.deepcopy(base_p._p)
                base_p._p.addnext(new_p)
                from pptx.text.text import _Paragraph

                wrapped = _Paragraph(new_p, tf)
                for r in list(wrapped.runs[1:]):
                    r._r.getparent().remove(r._r)
                wrapped.runs[0].text = line
                base_p = wrapped


def delete_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    rId = slides[index].get(qn("r:id"))
    prs.part.drop_rel(rId)
    xml_slides.remove(slides[index])


def build():
    prs = Presentation(str(SRC))

    icon_row_source = 6  # "你的 C 语言经验，已经懂了一半汇编"
    cmd_panel_source = 18  # "Makefile 的核心含义：别找错翻译官"

    concept_slide = duplicate_slide(prs, icon_row_source)
    set_texts(
        concept_slide,
        {
            "Text 0": "第 3 节",
            "Text 1": "Makefile 是什么：一份装配单",
            "Text 4": "目标: 依赖  写清楚要做出什么，需要先有什么才能做。",
            "Text 7": "命令必须 Tab 缩进  写成空格缩进会直接报错，这是新手最容易踩的坑。",
            "Text 10": ":= 和 ?=  := 是立刻赋值，?= 是只在变量还没被设置时才赋值。",
            "Text 13": ".PHONY  声明的目标不对应真实文件，每次执行都会重新跑一遍命令。",
        },
    )

    syntax_slide = duplicate_slide(prs, cmd_panel_source)
    set_texts(
        syntax_slide,
        {
            "Text 0": "第 3 节",
            "Text 1": "Makefile 语法速览（通用示例，不是本项目代码）",
            "Text 2": "先认识这份『装配单』的写法，再看下一页本项目 Makefile 怎么用它。",
            "Text 7": (
                "目标: 依赖1 依赖2\n"
                "\t命令              # 命令前必须是 Tab，不能是空格\n"
                "\n"
                "hello.o: hello.c\n"
                "\tgcc -c hello.c -o hello.o\n"
                "\n"
                "CC := gcc            # := 立即展开赋值\n"
                "CFLAGS ?= -Wall       # ?= 只在变量还没被设置时才赋值\n"
                "\n"
                ".PHONY: clean         # 声明伪目标，clean 不是一个真实文件\n"
                "clean:\n"
                "\trm -f *.o"
            ),
        },
    )

    # Keep only the two new slides; delete everything else, back to front
    # so indices of not-yet-deleted slides stay valid.
    total_original = len(prs.slides) - 2
    for i in range(total_original - 1, -1, -1):
        delete_slide(prs, i)

    prs.save(str(OUT))
    print(OUT)
    print("slides=", len(prs.slides))


if __name__ == "__main__":
    build()
