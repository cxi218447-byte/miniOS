"""Extend week01_qemu_hello_course_v2 (1).pptx with in-class lab slides.

Clones existing slide styles (command panel = slide 33, numbered list =
slide 36) via raw XML copy so the new slides match the deck's visual
design exactly, then fills in environment-check / troubleshooting /
acceptance-criteria content sourced from docs/week01/qemu_hello.md.
"""

import copy
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "docs" / "week01" / "week01_qemu_hello_course_v2 (1).pptx"
OUT = ROOT / "docs" / "week01" / "week01_qemu_hello_course.pptx"


def duplicate_slide(prs, index):
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)

    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    for shp in source.shapes:
        new_el = copy.deepcopy(shp._element)
        dest.shapes._spTree.append(new_el)

    for shp in dest.shapes:
        blip = shp._element.find(".//" + qn("a:blip"))
        if blip is not None:
            old_rid = blip.get(qn("r:embed"))
            if old_rid:
                image_part = source.part.related_part(old_rid)
                new_rid = dest.part.relate_to(
                    image_part,
                    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
                )
                blip.set(qn("r:embed"), new_rid)

    return dest


def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])


def find_slide_index(prs, exact_text):
    for i, slide in enumerate(prs.slides):
        for shp in slide.shapes:
            if shp.has_text_frame and shp.text_frame.text.strip() == exact_text:
                return i
    raise ValueError(f"no slide found with text: {exact_text!r}")


def set_texts(slide, text_by_shape_name):
    for shp in slide.shapes:
        if shp.name in text_by_shape_name and shp.has_text_frame:
            new_text = text_by_shape_name[shp.name]
            tf = shp.text_frame
            paragraphs = tf.paragraphs
            lines = new_text.split("\n")
            # reuse existing paragraph/run formatting, only swap text
            base_p = paragraphs[0]
            if not base_p.runs:
                continue
            base_run = base_p.runs[0]
            base_run.text = lines[0]
            for extra_run in list(base_p.runs[1:]):
                extra_run.text = ""
            # remove any extra pre-existing paragraphs beyond the first
            for p in list(paragraphs[1:]):
                p._p.getparent().remove(p._p)
            for line in lines[1:]:
                new_p = copy.deepcopy(base_p._p)
                base_p._p.addnext(new_p)
                from pptx.text.text import _Paragraph

                wrapped = _Paragraph(new_p, tf)
                for r in list(wrapped.runs[1:]):
                    r._r.getparent().remove(r._r)
                wrapped.runs[0].text = line
                base_p = wrapped


def insert_after(prs, new_slide_index, anchor_text):
    target = find_slide_index(prs, anchor_text) + 1
    move_slide(prs, new_slide_index, target)
    return target


def insert_before(prs, new_slide_index, anchor_text):
    target = find_slide_index(prs, anchor_text)
    move_slide(prs, new_slide_index, target)
    return target


def build():
    prs = Presentation(str(SRC))

    cmd_panel_source = find_slide_index(prs, "实验命令：三步跑通 Hello miniOS")
    numbered_list_source = find_slide_index(prs, "第一周学生作业")

    # --- New slide A: env check detail (command-panel style) ---
    env_slide = duplicate_slide(prs, cmd_panel_source)
    set_texts(
        env_slide,
        {
            "Text 0": "第 4 节",
            "Text 1": "上课第一步：环境检查怎么做",
            "Text 2": "缺工具不丢人，记录真实结果才重要——不要跳过这步直接编译。",
            "Text 7": (
                "# 需要确认这些命令存在\n"
                "make\n"
                "qemu-system-loongarch64\n"
                "loongarch64-linux-gnu-gcc\n"
                "loongarch64-linux-gnu-objdump\n"
                "\n"
                "# 一键检查（推荐先跑这个）\n"
                "sh scripts/check-env.sh"
            ),
        },
    )
    insert_before(prs, len(prs.slides) - 1, "实验命令：三步跑通 Hello miniOS")

    # --- New slide B: common errors (numbered-list style, 4 items) ---
    numbered_list_source = find_slide_index(prs, "第一周学生作业")
    error_slide = duplicate_slide(prs, numbered_list_source)
    set_texts(
        error_slide,
        {
            "Text 0": "第 4 节",
            "Text 1": "常见错误排查：卡住了先看这四条",
            "Text 4": "1",
            "Text 5": "command not found（make/gcc/qemu）→ 先跑 sh scripts/check-env.sh，按提示安装缺失工具。",
            "Text 8": "2",
            "Text 9": "make run 后没有 Hello 输出 → 先 make clean && make，确认生成了 build/minios.elf。",
            "Text 12": "3",
            "Text 13": "不知道怎么退出 QEMU → -nographic 模式下，先按 Ctrl-a，松开后再按 x。",
            "Text 16": "4",
            "Text 17": "终端里还看到 .data/.bss 相关输出 → 当前 master 已含第 2 周代码，第 1 周只关注 Hello 这一行。",
        },
    )
    insert_after(prs, len(prs.slides) - 1, "预期输出与测试记录")

    # --- New slide C: acceptance criteria (numbered-list style, 4 items) ---
    numbered_list_source = find_slide_index(prs, "第一周学生作业")
    accept_slide = duplicate_slide(prs, numbered_list_source)
    set_texts(
        accept_slide,
        {
            "Text 0": "第 4 节",
            "Text 1": "今天课上要做到这四条",
            "Text 4": "1",
            "Text 5": "环境检查 + make clean && make + make run，全部有真实记录，没跑就写「未执行」。",
            "Text 8": "2",
            "Text 9": "QEMU 终端里看到 Hello miniOS on LoongArch64。",
            "Text 12": "3",
            "Text 13": "能讲清楚：CPU 为什么先进 boot/start.S、为什么要设 $sp、为什么用 printk 不用 printf。",
            "Text 16": "4",
            "Text 17": "不强求：解释 .data/.bss 细节、修改启动代码、开发板迁移——这些是后面几周的内容。",
        },
    )
    insert_after(prs, len(prs.slides) - 1, "常见错误排查：卡住了先看这四条")

    prs.save(str(OUT))
    print(OUT)
    print("slides=", len(prs.slides))


if __name__ == "__main__":
    build()
