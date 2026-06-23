from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "week01_qemu_hello_4lessons_v4.pptx"
SOURCE_URL = "https://blog.csdn.net/loongsoner/article/details/128977015"

COLORS = {
    "ink": RGBColor(32, 35, 40),
    "muted": RGBColor(88, 96, 108),
    "bg": RGBColor(248, 249, 251),
    "line": RGBColor(212, 217, 226),
    "blue": RGBColor(31, 83, 158),
    "green": RGBColor(42, 125, 77),
    "orange": RGBColor(188, 100, 35),
    "teal": RGBColor(0, 113, 128),
    "red": RGBColor(168, 68, 62),
    "code_bg": RGBColor(34, 38, 46),
    "code_fg": RGBColor(238, 241, 246),
}


def font(run, size=22, bold=False, color=None, name="Microsoft YaHei"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["bg"]


def text(slide, x, y, w, h, body, size=22, bold=False, color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = body
    font(r, size=size, bold=bold, color=color)
    return box


def title(slide, body, subtitle=None, lesson=None):
    if lesson:
        text(slide, 0.7, 0.28, 2.0, 0.3, lesson, size=13, bold=True, color=COLORS["blue"])
    text(slide, 0.7, 0.55, 12.0, 0.55, body, size=30, bold=True)
    if subtitle:
        text(slide, 0.72, 1.12, 11.8, 0.32, subtitle, size=14, color=COLORS["muted"])
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(1.45), Inches(11.95), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.fill.background()


def bullets(prs, heading, items, subtitle=None, lesson=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, heading, subtitle, lesson)
    box = slide.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(11.25), Inches(5.1))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(9)
        r = p.add_run()
        r.text = item
        font(r, size=21 if len(item) < 35 else 18)
    return slide


def code(prs, heading, body, subtitle=None, lesson=None, size=15):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, heading, subtitle, lesson)
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.85),
        Inches(1.72),
        Inches(11.65),
        Inches(5.05),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["code_bg"]
    shape.line.color.rgb = COLORS["code_bg"]
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.22)
    tf.margin_top = Inches(0.18)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = body
    font(r, size=size, color=COLORS["code_fg"], name="Consolas")


def source_slide(prs, path, body, note, size=11):
    code(
        prs,
        f"项目代码：{path}",
        body,
        note,
        "第 4 节",
        size=size,
    )


def abi_table(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, "LoongArch 通用寄存器 ABI 别名表", "第 1 周先看懂别名体系，后续讲函数调用时再深入保存规则。", "第 2 节")
    rows = [
        ("r0", "zero", "常量 0", "硬件恒为 0"),
        ("r1", "ra", "返回地址", "bl 调用函数时写入"),
        ("r2", "tp", "线程指针", "TLS 相关，本周只认识"),
        ("r3", "sp", "栈指针", "进入 C 函数前必须设置"),
        ("r4-r5", "a0-a1 / v0-v1", "参数 / 返回值", "前两个参数，也可放返回值"),
        ("r6-r11", "a2-a7", "参数", "第 3 到第 8 个整型参数"),
        ("r12-r20", "t0-t8", "临时寄存器", "调用者保存，短期计算常用"),
        ("r21", "保留", "平台保留", "普通程序不要随意使用"),
        ("r22", "fp / s9", "帧指针 / 保存寄存器", "需要帧指针时作 fp"),
        ("r23-r31", "s0-s8", "保存寄存器", "被调用者保存"),
    ]
    table_shape = slide.shapes.add_table(len(rows) + 1, 4, Inches(0.75), Inches(1.8), Inches(11.85), Inches(4.8))
    table = table_shape.table
    widths = [1.35, 2.2, 2.35, 5.95]
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)
    headers = ["编号", "ABI 别名", "主要用途", "课堂说明"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["blue"]
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                font(r, size=13, bold=True, color=RGBColor(255, 255, 255))
    for row_idx, row in enumerate(rows, start=1):
        for col, value in enumerate(row):
            cell = table.cell(row_idx, col)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    font(r, size=11 if col == 3 else 12)
    text(slide, 0.78, 6.72, 11.9, 0.25, "记忆方式：先背 r0、r1、r3、r4-r11、r12-r20；r22-r31 到函数调用与栈帧实验再重点讲。", size=12, color=COLORS["muted"])


def section(prs, lesson, heading, goal):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    text(slide, 0.85, 1.15, 2.2, 0.38, lesson, size=15, bold=True, color=COLORS["blue"])
    text(slide, 0.85, 1.75, 11.4, 0.75, heading, size=36, bold=True)
    text(slide, 0.9, 2.85, 10.8, 0.9, goal, size=23, color=COLORS["muted"])
    text(slide, 0.9, 6.55, 11.3, 0.28, "本节结束时，学生要能用自己的话解释核心概念。", size=13, color=COLORS["muted"])


def cards(slide, card_data):
    x = 0.85
    for name, body, color in card_data:
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.25), Inches(2.25), Inches(1.7))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(255, 255, 255)
        rect.line.color.rgb = COLORS["line"]
        text(slide, x + 0.17, 2.43, 1.9, 0.3, name, size=17, bold=True, color=color)
        text(slide, x + 0.17, 2.87, 1.9, 0.85, body, size=14)
        x += 2.45


def flow(prs, heading, card_data, subtitle=None, lesson=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, heading, subtitle, lesson)
    cards(slide, card_data)
    x = 3.04
    for _ in range(len(card_data) - 1):
        text(slide, x, 2.85, 0.35, 0.28, "→", size=24, bold=True, color=COLORS["muted"], align=PP_ALIGN.CENTER)
        x += 2.45


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    text(slide, 0.85, 1.0, 11.8, 0.8, "第 1 周：从 C 语言走进 LoongArch 汇编", size=35, bold=True)
    text(slide, 0.9, 1.9, 11.3, 0.45, "4 节课导入 + QEMU Hello miniOS 实验", size=21, color=COLORS["blue"])
    text(slide, 0.9, 3.0, 11.5, 0.7, "对象：只系统学完 C 语言的大二计算机学生", size=24)
    text(slide, 0.9, 6.42, 11.6, 0.3, "参考：教材读书笔记《汇编语言编程基础 基于 LoongArch》；实验以本仓库 miniOS 代码为准。", size=13, color=COLORS["muted"])

    bullets(
        prs,
        "四节课安排",
        [
            "第 1 节：为什么学习汇编，机器语言、汇编语言、C 语言是什么关系",
            "第 2 节：LoongArch 最小知识包：寄存器、指令、寻址、分支",
            "第 3 节：从 C 到可执行文件：预处理、编译、汇编、链接、ELF",
            "第 4 节：读 miniOS 第 1 周代码，在 QEMU 中看到 Hello 输出",
        ],
        "第一周不追求掌握全部指令，而是建立正确的底层执行模型。",
    )

    section(prs, "第 1 节", "为什么学习汇编", "把学生熟悉的 C 语言，连接到 CPU 真正执行的指令。")
    bullets(
        prs,
        "第一节学习目标",
        [
            "能说清 C 语言、汇编语言、机器语言的层次关系",
            "知道指令集是软件和硬件之间的接口",
            "理解汇编语言和具体 CPU 架构绑定，不能随意跨架构运行",
            "知道本课程第一阶段只在 QEMU 中做 LoongArch64 实验",
        ],
        lesson="第 1 节",
    )
    flow(
        prs,
        "从 C 代码到 CPU 执行",
        [
            ("C 语言", "接近人的表达\n变量、函数、循环", COLORS["green"]),
            ("汇编语言", "接近 CPU 指令\n寄存器、跳转、访存", COLORS["blue"]),
            ("机器指令", "二进制编码\nCPU 直接取指执行", COLORS["orange"]),
            ("硬件行为", "运算、读写内存\n访问外设", COLORS["teal"]),
        ],
        "教材笔记强调：指令系统是软硬件之间的接口。",
        "第 1 节",
    )
    bullets(
        prs,
        "用 C 语言经验理解汇编",
        [
            "C 里的变量：很多时候会临时放在寄存器或内存里",
            "C 里的 if / while：底层会变成比较和跳转指令",
            "C 里的函数调用：底层涉及参数寄存器、返回地址和栈",
            "C 里的指针：底层就是地址，访存指令按地址读写数据",
        ],
        lesson="第 1 节",
    )
    bullets(
        prs,
        "第 1 节课堂活动",
        [
            "让学生写一个最简单的 C 函数：int add(int a, int b)",
            "提问：CPU 真的认识变量名 a 和 b 吗？",
            "提问：return a + b 最终需要哪些底层动作？",
            "结论：汇编课不是背指令，而是看清 C 程序如何落到机器执行",
        ],
        lesson="第 1 节",
    )

    section(prs, "第 2 节", "LoongArch 最小知识包", "只讲读懂第 1 周代码需要的寄存器、指令和寻址方式。")
    bullets(
        prs,
        "第二节学习目标",
        [
            "知道 LoongArch 指令长度固定为 32 位",
            "先认识通用寄存器编号 r0~r31，再认识 ABI 别名 zero、ra、sp、a0、t0",
            "理解三类基础指令：算术运算、访存、转移",
            "能读懂 boot/start.S 中的主路径",
        ],
        lesson="第 2 节",
    )
    bullets(
        prs,
        "寄存器：先记 r 编号，再记 ABI 别名",
        [
            "r0 / zero：常量寄存器，值永远是 0",
            "r1 / ra：函数返回地址，bl 调用函数时会写入返回地址",
            "r3 / sp：栈指针，C 函数运行前必须准备好",
            "r4~r11 / a0~a7：参数寄存器，其中 r4~r5 也可作返回值 v0~v1",
            "r12~r20 / t0~t8：临时寄存器，适合短期计算",
        ],
        lesson="第 2 节",
    )
    abi_table(prs)
    code(
        prs,
        "教材笔记风格的基础指令例子",
        """addi.d  r3, r3, -32       # 64 位加法：sp = sp - 32
st.d    r1, r3, 24        # 把 ra 保存到 sp + 24
ld.d    r1, r3, 24        # 从 sp + 24 恢复 ra
addi.w  r8, r1, 16        # 32 位加法，立即数参与运算
bne     r12, r0, L        # r12 != 0 时跳转到 L
bl      func              # 调用 func，返回地址写入 r1/ra
jirl    r0, r1, 0         # 跳回 r1/ra 指向的位置

# 教材笔记还提醒：move 是宏指令，可由 or rd,rj,r0 表达。""",
        "课堂先用 r0~r31 讲清楚，再补充 zero/ra/sp 等 ABI 别名。",
        "第 2 节",
    )
    code(
        prs,
        "本仓库代码里的写法和教材写法如何对应",
        """# 教材/反汇编里常见：
addi.d  r3, r3, -32
st.d    r1, r3, 24
jirl    r0, r1, 0

# GNU 汇编器也接受 ABI 别名：
addi.d  $sp, $sp, -32     # $sp 就是 r3
st.d    $ra, $sp, 24      # $ra 就是 r1
jirl    $zero, $ra, 0     # $zero 就是 r0

# la.global 这类用于装入符号地址，按伪指令/宏指令讲，
# 不把它当作 LoongArch 基础整数指令本体。""",
        "这页专门避免学生误以为课件和教材是两套汇编。",
        "第 2 节",
    )
    bullets(
        prs,
        "寻址方式用一句话理解",
        [
            "寄存器寻址：如 addi.d r3, r3, -32，操作数来自寄存器",
            "立即数寻址：如 addi.w r8, r1, 16，常数 16 写在指令中",
            "基址 + 偏移寻址：如 ld.d r1, r3, 24，访问 r3 + 24 指向的内存",
            "PC 相对转移：如 bl、bne，跳转目标由当前指令位置和偏移计算",
        ],
        "教材笔记把寻址方式作为理解访存和跳转的关键入口。",
        "第 2 节",
    )
    bullets(
        prs,
        "第二节练习",
        [
            "指出 addi.d、ld.d、st.d、bne 分别属于哪类指令",
            "把 bne r12, r0, L 解释成 C 语言里的 if 条件跳转",
            "解释 bl 为什么会影响 r1/ra",
            "解释为什么设置 r3/sp 后才能放心进入 C 函数",
        ],
        lesson="第 2 节",
    )

    section(prs, "第 3 节", "从 C 到可执行文件", "让学生知道 Makefile 背后发生了什么。")
    bullets(
        prs,
        "第三节学习目标",
        [
            "理解 .c、.S、.o、.elf 文件各自是什么",
            "知道编译流程大致包括预处理、编译、汇编、链接",
            "知道 .S 文件会先经过 C 预处理器",
            "理解链接脚本负责规定入口和内存布局",
        ],
        lesson="第 3 节",
    )
    flow(
        prs,
        "教材笔记中的编译流程，用到本实验里",
        [
            (".c/.S", "源代码\nC 和汇编", COLORS["green"]),
            ("预处理/编译", "生成汇编\n处理 include", COLORS["blue"]),
            ("汇编", "生成 .o\n机器指令片段", COLORS["orange"]),
            ("链接", "合成 .elf\n确定地址和入口", COLORS["teal"]),
            ("QEMU", "加载内核\n开始执行", COLORS["red"]),
        ],
        "第 1 周学生只需要能说出每一步的作用。",
        "第 3 节",
    )
    code(
        prs,
        "本实验 Makefile 的核心含义",
        """CROSS_COMPILE ?= loongarch64-linux-gnu-
CC      := $(CROSS_COMPILE)gcc
OBJCOPY := $(CROSS_COMPILE)objcopy
QEMU    ?= qemu-system-loongarch64

TARGET  := build/minios.elf
LDFLAGS := -T kernel/linker.ld -nostdlib -static""",
        "强调：不能用宿主机 x86_64 gcc 冒充 LoongArch 工具链。",
        "第 3 节",
    )
    code(
        prs,
        "链接脚本只先讲三件事",
        """ENTRY(_start)

SECTIONS
{
    . = 0x9000000000200000;

    .text : { *(.text.boot) *(.text*) }
    .rodata : { *(.rodata*) }
    .data : { *(.data*) }
    .bss : { *(.bss*) *(COMMON) }
}""",
        "入口是谁、从哪个地址开始、各类内容放在哪里。",
        "第 3 节",
    )
    bullets(
        prs,
        "普通程序和裸机程序的链接差异",
        [
            "普通 Linux 程序：C 运行库和操作系统会准备运行环境",
            "miniOS 裸机程序：没有 libc，没有操作系统帮忙调用 main",
            "所以我们需要 _start、栈、链接脚本、串口输出",
            "这正是第 1 周 Hello miniOS 的教学价值",
        ],
        lesson="第 3 节",
    )
    bullets(
        prs,
        "第三节课堂活动",
        [
            "在黑板画出 .c/.S → .o → .elf → QEMU 的流程",
            "让学生找出 Makefile 中 CC、TARGET、LDFLAGS、run 四个位置",
            "解释为什么 .elf 不是 Windows 可直接双击运行的程序",
            "提醒：本周不深入 ELF 格式，后续再展开",
        ],
        lesson="第 3 节",
    )

    section(prs, "第 4 节", "QEMU Hello miniOS 实验", "把前 3 节概念落到仓库代码和真实运行命令。")
    bullets(
        prs,
        "第四节学习目标",
        [
            "能指出第 1 周核心代码文件的位置",
            "能按执行路径读懂 _start 到 kernel_main",
            "能解释 printk 为什么能把字符显示到终端",
            "能完成环境检查、编译、运行和测试记录",
        ],
        lesson="第 4 节",
    )
    bullets(
        prs,
        "第 1 周核心文件",
        [
            "boot/start.S：启动入口，设置栈，跳转 kernel_main",
            "kernel/linker.ld：规定入口地址和段布局",
            "kernel/main.c：C 语言内核入口，调用 printk",
            "kernel/printk.c：最小字符串输出",
            "include/printk.h、include/uart.h、include/types.h：头文件",
            "Makefile：编译、运行和清理命令",
        ],
        lesson="第 4 节",
    )
    bullets(
        prs,
        "先说明：当前仓库代码比第 1 周略多",
        [
            "boot/start.S 中已经包含 .bss 清零和 exception_entry，课堂第 1 周先讲 _start 主路径",
            "kernel/main.c 中已经包含 data/bss/string 测试，课堂第 1 周先抓 printk 输出链路",
            "lib/string.S、kernel/exception.c、kernel/syscall.c 属于后续周次铺垫，第一周不展开",
            "PPT 后续代码页会标注“第一周关注点”，防止学生一开始被扩展内容淹没",
        ],
        lesson="第 4 节",
    )
    source_slide(
        prs,
        "boot/start.S",
        """    .section .text.boot, "ax"
    .globl _start

_start:
    /* 设置内核栈。栈从高地址向低地址增长。 */
    la.global   $sp, boot_stack_top

    /* 清零 .bss，保证未初始化的全局变量从 0 开始。 */
    bl          clear_bss

    /* 进入 C 语言内核主函数。 */
    bl          kernel_main

halt:
    idle        0
    b           halt""",
        "第一周关注：_start、设置 sp、调用 kernel_main、停在 halt 循环。",
        size=12,
    )
    source_slide(
        prs,
        "boot/start.S 中的 clear_bss",
        """clear_bss:
    la.global   $t0, __bss_start
    la.global   $t1, __bss_end

1:
    beq         $t0, $t1, 2f
    st.b        $zero, $t0, 0
    addi.d      $t0, $t0, 1
    b           1b

2:
    jr          $ra""",
        "这段是第 2 周 .bss 的铺垫。第一周只需知道它在进入 C 前做内存初始化。",
        size=13,
    )
    source_slide(
        prs,
        "kernel/linker.ld",
        """ENTRY(_start)

SECTIONS
{
    . = 0x9000000000200000;

    .text : ALIGN(4K) {
        KEEP(*(.text.boot))
        *(.text .text.*)
    }

    .rodata : ALIGN(4K) { *(.rodata .rodata.*) }
    .data   : ALIGN(4K) { *(.data .data.*) }

    __bss_start = .;
    .bss : ALIGN(4K) {
        *(.bss .bss.*)
        *(COMMON)
    }
    __bss_end = .;
}""",
        "第一周关注：入口是 _start，.text.boot 被保留并放在前面，链接地址面向 QEMU virt。",
        size=11,
    )
    source_slide(
        prs,
        "kernel/main.c",
        """#include "printk.h"
#include "string.h"

static char bss_buffer[16];
static char data_message[] = "data section ok";

void kernel_main(void)
{
    char buf[32];
    const char *msg = "Hello, LoongArch miniOS!\\n";

    printk("miniOS booting...\\n");
    printk(msg);

    memset(buf, 0, sizeof(buf));
    memcpy(buf, data_message, strlen(data_message));
    printk(buf);
    printk("\\n");

    if (bss_buffer[0] == 0) {
        printk("bss section cleared\\n");
    }

    printk("week1-week2 check done\\n");
    while (1) { __asm__ volatile("idle 0"); }
}""",
        "当前文件已包含第 2/4 周铺垫。第一周课堂先只追踪 printk(msg) 到 UART 的路径。",
        size=10,
    )
    source_slide(
        prs,
        "kernel/printk.c",
        """#include "printk.h"
#include "uart.h"

void uart_putc(char ch)
{
    volatile unsigned char *uart =
        (volatile unsigned char *)UART0_BASE;

    /* 裸机早期先不轮询状态寄存器，直接写发送寄存器。 */
    *uart = (unsigned char)ch;
}

void uart_puts(const char *s)
{
    while (*s) {
        if (*s == '\\n') {
            uart_putc('\\r');
        }
        uart_putc(*s++);
    }
}

void printk(const char *s)
{
    uart_puts(s);
}""",
        "第一周关注：字符串逐字符输出，最终是向 UART0_BASE 指向的 MMIO 地址写字节。",
        size=10,
    )
    source_slide(
        prs,
        "include/printk.h 与 include/uart.h",
        """/* include/printk.h */
#ifndef MINIOS_PRINTK_H
#define MINIOS_PRINTK_H

void printk(const char *s);
void printk_hex(unsigned long value);

#endif

/* include/uart.h */
#ifndef MINIOS_UART_H
#define MINIOS_UART_H

#include "types.h"

#define UART0_BASE 0x900000001fe001e0UL

void uart_putc(char ch);
void uart_puts(const char *s);

#endif""",
        "头文件告诉 C 编译器函数原型和 UART 地址。UART0_BASE 只用于 QEMU virt，不代表 2K0300。",
        size=10,
    )
    source_slide(
        prs,
        "include/types.h 与 include/string.h",
        """/* include/types.h */
#ifndef MINIOS_TYPES_H
#define MINIOS_TYPES_H

typedef unsigned long size_t;
typedef unsigned long uint64_t;
typedef long int64_t;

#endif

/* include/string.h */
#ifndef MINIOS_STRING_H
#define MINIOS_STRING_H

#include "types.h"

void *memset(void *dst, int value, size_t n);
void *memcpy(void *dst, const void *src, size_t n);
size_t strlen(const char *s);

#endif""",
        "这些是当前 main.c 需要的声明。string 函数属于后续周次，第一周不要求实现细节。",
        size=10,
    )
    source_slide(
        prs,
        "Makefile：工具链与目标文件",
        """CROSS_COMPILE ?= loongarch64-linux-gnu-
CC      := $(CROSS_COMPILE)gcc
OBJCOPY := $(CROSS_COMPILE)objcopy
QEMU    ?= qemu-system-loongarch64

BUILD_DIR := build
TARGET    := $(BUILD_DIR)/minios.elf
BIN       := $(BUILD_DIR)/minios.bin

CFLAGS  := -Wall -Wextra -O2 -g -ffreestanding
CFLAGS  += -fno-builtin -fno-stack-protector
CFLAGS  += -nostdlib -mabi=lp64d -march=loongarch64

LDFLAGS := -T kernel/linker.ld -nostdlib -static""",
        "第一周关注：使用 LoongArch64 交叉编译器，生成裸机 ELF，不链接宿主机 C 运行库。",
        size=10,
    )
    source_slide(
        prs,
        "Makefile：编译、运行、清理",
        """all: $(TARGET) $(BIN)

$(TARGET): $(OBJS) kernel/linker.ld
    $(CC) $(CFLAGS) $(LDFLAGS) -o $@ $(OBJS)

$(BIN): $(TARGET)
    $(OBJCOPY) -O binary $< $@

$(BUILD_DIR)/%.o: %.c
    mkdir -p $(dir $@)
    $(CC) $(CFLAGS) -Iinclude -c $< -o $@

$(BUILD_DIR)/%.o: %.S
    mkdir -p $(dir $@)
    $(CC) $(ASFLAGS) -Iinclude -c $< -o $@

run: $(TARGET)
    $(QEMU) $(QEMU_ARGS)

clean:
    rm -rf $(BUILD_DIR)""",
        "课堂演示时让学生对应 make、make run、make clean 三个动作。",
        size=10,
    )
    flow(
        prs,
        "第 1 周最小运行链路",
        [
            ("make", "交叉编译\n生成 ELF", COLORS["blue"]),
            ("QEMU", "模拟 virt\n加载内核", COLORS["teal"]),
            ("_start", "设置 sp\n调用 C", COLORS["orange"]),
            ("kernel_main", "调用 printk\n输出字符串", COLORS["green"]),
            ("UART", "终端显示\nHello", COLORS["red"]),
        ],
        "这条链路必须跑通，才进入后续实验。",
        "第 4 节",
    )
    code(
        prs,
        "boot/start.S 主路径：先用教材风格读",
        """    .section .text.boot, "ax"
    .globl _start

_start:
    # 设置内核栈：r3/sp 指向 boot_stack_top
    la.global   r3, boot_stack_top

    # 进入 C 语言内核主函数
    bl          kernel_main

halt:
    idle        0
    b           halt""",
        "la.global 是装入符号地址的汇编器伪指令；核心动作是让 r3/sp 成为可用栈指针。",
        "第 4 节",
    )
    code(
        prs,
        "同一段代码：源码里的 ABI 别名写法",
        """_start:
    # $sp 是 r3 的 ABI 别名
    la.global   $sp, boot_stack_top

    # bl 调用函数，返回地址写入 r1/ra
    bl          kernel_main

halt:
    idle        0
    b           halt""",
        "课堂讲解顺序：先 r 编号，再 ABI 别名；不要一上来只讲 $sp。",
        "第 4 节",
    )
    code(
        prs,
        "建议第 1 周使用的最小 kernel_main",
        """#include "printk.h"

void kernel_main(void)
{
    printk("Hello miniOS on LoongArch64\\n");

    while (1) {
        __asm__ volatile("idle 0");
    }
}""",
        "当前仓库已有后续扩展痕迹，课堂讲解时先抽出这条最小路径。",
        "第 4 节",
    )
    code(
        prs,
        "UART 输出：C 代码背后是访存",
        """void uart_putc(char ch)
{
    volatile unsigned char *uart =
        (volatile unsigned char *)UART0_BASE;

    *uart = (unsigned char)ch;
}

void printk(const char *s)
{
    while (*s) {
        uart_putc(*s++);
    }
}""",
        "编译后会落到地址计算和 st.b/st.w 一类访存指令；UART0_BASE 是 QEMU virt 平台地址。",
        "第 4 节",
    )
    code(
        prs,
        "实验命令",
        """# 环境检查
which make
which qemu-system-loongarch64
which loongarch64-linux-gnu-gcc

# 编译
make clean
make

# 运行
make run""",
        "实际测试结果必须来自学生机器真实命令输出。",
        "第 4 节",
    )
    bullets(
        prs,
        "预期输出与测试记录",
        [
            "预期串口输出：Hello miniOS on LoongArch64",
            "未执行过 make run，就不能写“已通过”",
            "工具链缺失时，记录为未执行，并写明失败原因",
            "测试报告必须包含：已执行、未执行、失败原因、下一步命令",
        ],
        lesson="第 4 节",
    )
    bullets(
        prs,
        "常见错误和排查",
        [
            "找不到 loongarch64-linux-gnu-gcc：安装交叉编译器",
            "找不到 qemu-system-loongarch64：安装 qemu-system-misc",
            "误用系统 gcc：检查 Makefile 中 CROSS_COMPILE",
            "QEMU 无输出：检查入口、链接地址、UART 地址、run 命令",
            "中文显示异常：先确认文件编码和终端显示编码，不急着改源码",
        ],
        lesson="第 4 节",
    )
    bullets(
        prs,
        "第一周学生作业",
        [
            "画出 miniOS 第 1 周执行路径图",
            "把 Hello 字符串改成自己的姓名和学号，重新编译运行",
            "解释 boot/start.S 中 sp、bl、b、idle 的作用",
            "提交环境检查结果和真实运行截图或日志",
        ],
        lesson="第 4 节",
    )
    bullets(
        prs,
        "教师板书建议",
        [
            "C 语言：函数、变量、循环",
            "汇编语言：寄存器、指令、跳转、访存",
            "裸机启动：没有操作系统，所以要自己准备入口和栈",
            "实验闭环：改代码 → 编译 → QEMU 运行 → 记录结果",
        ],
        lesson="第 4 节",
    )
    bullets(
        prs,
        "本周边界",
        [
            "不讲完整操作系统",
            "不讲进程、文件系统、虚拟内存",
            "不直接适配 2K0300 开发板",
            "不把未实测结果写成已通过",
            "下一周再展开 .data、.bss 和内存初始化",
        ],
        lesson="第 4 节",
    )
    bullets(
        prs,
        "参考资料",
        [
            f"教材读书笔记：{SOURCE_URL}",
            "龙芯架构参考手册：用于查证指令、寄存器和特权架构细节",
            "LoongArch ELF psABI：用于后续函数调用约定、栈布局和 ABI 讲解",
            "本仓库代码和 docs/environment_check.md：用于实验实测记录",
        ],
    )

    try:
        prs.save(OUT)
        saved = OUT
    except PermissionError:
        saved = OUT.with_name(f"{OUT.stem}_4lessons{OUT.suffix}")
        prs.save(saved)
    print(saved)
    print(f"slides={len(prs.slides)}")


if __name__ == "__main__":
    build()
